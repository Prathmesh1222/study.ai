from pptx import Presentation
from pathlib import Path

# -------- Paths --------
BASE_DIR = Path(__file__).resolve().parent.parent
INPUT_DIR = BASE_DIR / "data" / "raw" / "ppt"
OUTPUT_DIR = BASE_DIR / "data" / "extracted_text"

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def extract_text_from_ppt(ppt_path):
    """
    Extracts slide-wise text from a PPT file
    """
    prs = Presentation(ppt_path)
    content = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        content.append(f"\n[SLIDE {slide_no}]")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        content.append(text)

    return "\n".join(content)


def process_all_ppts():
    """
    Processes ALL .pptx files inside data/raw/ppt (recursive)
    """
    ppt_files = list(INPUT_DIR.rglob("*.pptx"))

    if not ppt_files:
        print("❌ No PPT files found in data/raw/ppt/")
        return

    for ppt_file in ppt_files:
        unit_name = ppt_file.parent.name

        print(f"📄 Processing: {ppt_file.name}")

        header = (
            f"[UNIT: {unit_name}]\n"
            f"[SOURCE_FILE: {ppt_file.name}]\n"
            f"{'-'*40}\n"
        )

        extracted_text = extract_text_from_ppt(ppt_file)
        final_text = header + extracted_text

        output_file = OUTPUT_DIR / f"{ppt_file.stem}.txt"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(final_text)

        print(f"✅ Saved: {output_file}\n")


if __name__ == "__main__":
    process_all_ppts()
    